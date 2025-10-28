"""
Master Pipeline - Complete end-to-end research and validation agent calls
Usage: python pipeline.py <pitch_deck_file>

This script orchestrates:
1. Orchestrator Agent → Analyzes deck and creates research plan
2. Research Agent → Executes research tasks
3. Validator Agent → Validates claims
4. Generates final investment report
"""

import os
import sys
import json
import time
import subprocess
from pathlib import Path
from dotenv import load_dotenv

load_dotenv()


def run_command(cmd: list, description: str) -> bool:
    """Run a command and handle errors"""
    print(f"\n{'='*70}")
    print(f"▶️  {description}")
    print(f"{'='*70}")
    
    try:
        result = subprocess.run(cmd, check=True, capture_output=False)
        return result.returncode == 0
    except subprocess.CalledProcessError as e:
        print(f"\n❌ Failed: {description}")
        print(f"Error: {e}")
        return False


def extract_deck_text(deck_file: Path) -> str:
    """Extract text from pitch deck (PDF or PPT)"""
    print(f"\n📄 Extracting text from: {deck_file.name}")
    
    suffix = deck_file.suffix.lower()
    
    if suffix == '.pdf':
        try:
            import PyPDF2
            with open(deck_file, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                print(f"✅ Extracted {len(text)} characters from PDF")
                return text
        except Exception as e:
            print(f"⚠️  PDF extraction failed: {e}")
            print("   Trying alternative method...")
            
            try:
                import pdfplumber
                with pdfplumber.open(deck_file) as pdf:
                    text = ""
                    for page in pdf.pages:
                        text += page.extract_text() + "\n"
                    print(f"✅ Extracted {len(text)} characters with pdfplumber")
                    return text
            except Exception as e2:
                print(f"❌ Alternative PDF method also failed: {e2}")
                return ""
    
    elif suffix in ['.ppt', '.pptx']:
        try:
            from pptx import Presentation
            prs = Presentation(deck_file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            print(f"✅ Extracted {len(text)} characters from PowerPoint")
            return text
        except Exception as e:
            print(f"❌ PowerPoint extraction failed: {e}")
            return ""
    
    elif suffix in ['.txt', '.md']:
        try:
            with open(deck_file, 'r', encoding='utf-8') as f:
                text = f.read()
            print(f"✅ Read {len(text)} characters from text file")
            return text
        except Exception as e:
            print(f"❌ Text file read failed: {e}")
            return ""
    
    else:
        print(f"❌ Unsupported file type: {suffix}")
        return ""


def main():
    print("="*70)
    print("🚀 MASTER PIPELINE - Complete Due Diligence System")
    print("="*70)
    
    # Check API key
    api_key = os.environ.get("OPENROUTER_API_KEY")
    if not api_key:
        print("\n❌ Error: OPENROUTER_API_KEY not set")
        print("   Please set your OpenRouter API key:")
        print("   export OPENROUTER_API_KEY='your-key-here'")
        sys.exit(1)
    
    # Check for deck file
    if len(sys.argv) < 2:
        print("\n📖 Usage:")
        print("   python master_pipeline.py <pitch_deck_file>")
        print("\nSupported formats: .pdf, .pptx, .ppt, .txt, .md")
        print("\nExample:")
        print("   python master_pipeline.py startup_deck.pdf")
        sys.exit(1)
    
    deck_file = Path(sys.argv[1])
    
    if not deck_file.exists():
        print(f"\n❌ File not found: {deck_file}")
        sys.exit(1)
    
    print(f"\n📁 Input: {deck_file.name}")
    
    # Create output directory
    output_dir = Path("output")
    output_dir.mkdir(exist_ok=True)
    
    base_name = deck_file.stem
    
    # Step 1: Extract deck text
    print("\n" + "="*70)
    print("STEP 1: EXTRACT DECK CONTENT")
    print("="*70)
    
    deck_text = extract_deck_text(deck_file)
    
    if not deck_text or len(deck_text) < 100:
        print("\n❌ Failed to extract meaningful text from deck")
        print("   Please provide a deck with readable text content")
        sys.exit(1)
    
    # Save extracted text for reference
    text_file = output_dir / f"{base_name}_extracted_text.txt"
    with open(text_file, 'w', encoding='utf-8') as f:
        f.write(deck_text)
    print(f"✅ Saved extracted text: {text_file}")
    
    # Step 2: Run Orchestrator
    print("\n" + "="*70)
    print("STEP 2: ORCHESTRATOR ANALYSIS")
    print("="*70)
    
    from orchestrator import OrchestratorAgent
    
    orchestrator = OrchestratorAgent(api_key)
    research_plan = orchestrator.orchestrate(deck_text)
    
    if not research_plan:
        print("❌ Orchestrator failed to create research plan")
        sys.exit(1)
    
    # Save deck analysis
    deck_analysis_file = output_dir / f"{base_name}_deck_analysis.json"
    deck_analysis = orchestrator.analyze_deck(deck_text)
    with open(deck_analysis_file, 'w', encoding='utf-8') as f:
        json.dump(deck_analysis, f, indent=2, ensure_ascii=False)
    print(f"✅ Saved deck analysis: {deck_analysis_file}")
    
    # Save research plan
    research_plan_file = output_dir / f"{base_name}_research_plan.json"
    with open(research_plan_file, 'w', encoding='utf-8') as f:
        json.dump(research_plan.to_dict(), f, indent=2, ensure_ascii=False)
    print(f"✅ Saved research plan: {research_plan_file}")
    
    # Step 3: Execute Research
    print("\n" + "="*70)
    print("STEP 3: RESEARCH EXECUTION")
    print("="*70)
    
    print(f"\n⚠️  About to execute {len(research_plan.tasks)} research tasks")
    print(f"💰 Estimated cost: ${research_plan.estimated_cost_usd:.2f}")
    print(f"⏱️  Estimated time: {len(research_plan.tasks) * 30} seconds")
    
    confirm = input("\nProceed with research? (y/n): ")
    if confirm.lower() != 'y':
        print("❌ Pipeline cancelled by user")
        sys.exit(0)
    
    # Run research using subprocess to use the existing run_research.py
    research_cmd = [sys.executable, "run_research.py", str(research_plan_file)]
    
    if not run_command(research_cmd, "Executing Research Tasks"):
        print("⚠️  Research execution had issues, but continuing...")
    
    # Check for research results
    research_results_file = output_dir / f"{base_name}_research_results.json"
    validation_plan_file = output_dir / f"{base_name}_validation_plan.json"
    
    if not validation_plan_file.exists():
        print("❌ Validation plan not created by research agent")
        sys.exit(1)
    
    # Step 4: Execute Validation
    print("\n" + "="*70)
    print("STEP 4: VALIDATION EXECUTION")
    print("="*70)
    
    with open(validation_plan_file, 'r', encoding='utf-8') as f:
        validation_plan = json.load(f)
    
    print(f"\n⚠️  About to validate {validation_plan.get('total_tasks', 0)} claims")
    print(f"💰 Estimated cost: ${validation_plan.get('total_tasks', 0) * 0.03:.2f}")
    
    confirm = input("\nProceed with validation? (y/n): ")
    if confirm.lower() != 'y':
        print("❌ Pipeline cancelled by user")
        sys.exit(0)
    
    # Run validation
    validation_cmd = [sys.executable, "run_validation.py", str(validation_plan_file)]
    
    if not run_command(validation_cmd, "Executing Validation Tasks"):
        print("⚠️  Validation had issues, but checking for results...")
    
    # Step 5: Summary
    print("\n" + "="*70)
    print("✅ PIPELINE COMPLETE")
    print("="*70)
    
    # Check for final outputs
    final_report = output_dir / f"{base_name}_FINAL_REPORT.md"
    
    if final_report.exists():
        print(f"\n📊 FINAL REPORT GENERATED")
        print(f"   📄 {final_report}")
        
        # Show quick preview
        try:
            with open(final_report, 'r', encoding='utf-8') as f:
                lines = f.readlines()
                # Find the recommendation line
                for line in lines[:30]:
                    if "Investment Recommendation" in line or "RECOMMENDATION" in line:
                        print(f"\n🎯 {line.strip()}")
                        # Get next few lines
                        idx = lines.index(line)
                        for l in lines[idx+1:idx+5]:
                            if l.strip():
                                print(f"   {l.strip()}")
                        break
        except Exception as e:
            print(f"⚠️  Could not preview report: {e}")
    
    print(f"\n📁 All outputs saved to: {output_dir}/")
    print(f"\n🎉 Due diligence complete! Review {final_report.name} for details.")
    
    return True


if __name__ == "__main__":
    try:
        start_time = time.time()
        success = main()
        elapsed = time.time() - start_time
        
        print(f"\n⏱️  Total pipeline time: {elapsed/60:.1f} minutes")
        sys.exit(0 if success else 1)
        
    except KeyboardInterrupt:
        print("\n\n⚠️  Pipeline interrupted by user")
        sys.exit(1)
    except Exception as e:
        print(f"\n❌ Pipeline error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)